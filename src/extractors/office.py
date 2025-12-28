import docx
import pptx
import pandas as pd
import os
from .base import BaseExtractor

class DocxExtractor(BaseExtractor):
    def extract(self, file_path):
        try:
            doc = docx.Document(file_path)
            full_text = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        full_text.append(cell.text)
            text = "\n".join(full_text)
            if text.strip(): yield 1, text
        except Exception as e:
            print(f"⚠️ Docx Error {file_path}: {e}")

class SlideExtractor(BaseExtractor):
    def extract(self, file_path):
        try:
            prs = pptx.Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
                content = "\n".join(text_runs)
                if content.strip(): yield i + 1, content
        except Exception:
            pass

class SpreadsheetExtractor(BaseExtractor):
    def extract(self, file_path):
        try:
            # Check extension to decide method
            ext = os.path.splitext(file_path)[1].lower()
            
            if ext == '.csv':
                # Handle CSV
                df = pd.read_csv(file_path, nrows=20)
                text = df.to_string(index=False)
                if text.strip(): yield 1, f"CSV Data:\n{text}"
            else:
                # Handle Excel (.xls, .xlsx)
                xls = pd.ExcelFile(file_path)
                for sheet_name in xls.sheet_names:
                    df = pd.read_excel(xls, sheet_name=sheet_name, nrows=20)
                    text = df.to_string(index=False)
                    if text.strip(): yield 1, f"Sheet: {sheet_name}\n{text}"
                    
        except Exception as e:
            print(f"⚠️ Spreadsheet Error {file_path}: {e}")

class TextExtractor(BaseExtractor):
    def extract(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
                if text.strip(): yield 1, text
        except Exception:
            pass