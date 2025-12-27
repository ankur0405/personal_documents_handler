"""
Module: Content Extractor (Email Edition)
Description: Extracts text from PDFs, DOCX, PPTX, Images, and OUTLOOK MSG files.
"""

import pandas as pd
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import extract_msg  # <--- NEW LIBRARY
import easyocr
import os
from src.common.image_classifier import is_scanned_document

# Initialize OCR (Global)
print("👁️  Initializing Optical Character Recognition (OCR)...")
ocr_reader = easyocr.Reader(['en'], gpu=True) 

def extract_pages(file_path: str, file_type: str):
    """
    Generator that yields tuples: (page_number, text_content)
    """
    try:
        if file_type == 'pdf':
            yield from _read_pdf_pages(file_path)
        elif file_type == 'docx':
            yield from _read_docx_pages(file_path)
        elif file_type == 'pptx':
            yield from _read_pptx_pages(file_path)
            
        # --- NEW: EMAIL HANDLING ---
        elif file_type == 'msg':
            # Emails don't have "pages", so we return everything as Page 1.
            text = _read_msg_content(file_path)
            if text:
                yield (1, text)

        elif file_type in ['txt', 'md', 'csv', 'xml', 'html', 'json']:
            yield (1, _read_text(file_path))
            
        elif file_type in ['png', 'jpg', 'jpeg', 'heic']:
            if is_scanned_document(file_path):
                print(f"   📸 OCR Scanning: {os.path.basename(file_path)}")
                text = _read_image_ocr(file_path)
                if text:
                    yield (1, text)
            
        # --- SPREADSHEETS (NEW LOGIC) ---
        elif file_type in ['xlsx', 'xls', 'csv']:
            # We treat the summary of the sheet as "Page 1"
            text = _read_spreadsheet_summary(file_path)
            if text:
                yield (1, text)
            
    except Exception as e:
        print(f"⚠️  Skipping content for {os.path.basename(file_path)}: {e}")
        return

def _read_msg_content(path: str) -> str:
    """
    Parses Outlook .msg files.
    Extracts: Subject, Sender, Date, and Body.
    """
    try:
        msg = extract_msg.Message(path)
        
        # We construct a structured string so the AI understands the context
        content = (
            f"Subject: {msg.subject}\n"
            f"From: {msg.sender}\n"
            f"Date: {msg.date}\n"
            f"Body:\n{msg.body}"
        )
        msg.close() # Good practice to close the file handle
        return content
    except Exception as e:
        print(f"⚠️  Error parsing email {path}: {e}")
        return ""

# ... (Keep _read_image_ocr, _read_pdf_pages, _read_docx_pages, _read_pptx_pages, _read_text as they were) ...
# (If you need the full file again, let me know, but you just need to append the _read_msg_content function 
#  and update the extractor logic above).

def _read_image_ocr(path: str) -> str:
    try:
        results = ocr_reader.readtext(path, detail=0)
        return " ".join(results)
    except Exception:
        return ""

def _read_pdf_pages(path: str):
    with fitz.open(path) as doc:
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()
            if text.strip():
                yield (page_num, text)

def _read_docx_pages(path: str):
    doc = docx.Document(path)
    full_text = "\n".join([para.text for para in doc.paragraphs])
    if full_text.strip():
        yield (1, full_text)

def _read_pptx_pages(path: str):
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        full_text = "\n".join(text)
        if full_text.strip():
            yield (i, full_text)

def _read_text(path: str) -> str:
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def _read_spreadsheet_summary(path: str) -> str:
    """
    Reads the First 5 rows (Headers) and Last 5 rows (Totals) of a spreadsheet.
    """
    try:
        # Read the file into a DataFrame
        if path.endswith('.csv'):
            df = pd.read_csv(path)
        else:
            df = pd.read_excel(path)
            
        # 1. Clean up: Drop completely empty rows/cols
        df.dropna(how='all', inplace=True)
        
        # 2. Grab the HEAD (Headers + Context)
        head_text = df.head(5).to_string(index=False)
        
        # 3. Grab the TAIL (Totals usually live here)
        # We check if the dataframe is long enough to have a separate tail
        if len(df) > 10:
            tail_text = df.tail(5).to_string(index=False, header=False) # No header for tail
            summary = f"--- START OF SHEET ---\n{head_text}\n...\n[...Middle Rows Skipped...]\n...\n--- END OF SHEET (TOTALS) ---\n{tail_text}"
        else:
            # If file is small, just take the whole thing
            summary = df.to_string(index=False)
            
        return summary

    except Exception as e:
        print(f"⚠️  Error reading spreadsheet {path}: {e}")
        return ""