"""
Module: Extractors
Description: Concrete implementations of the BaseExtractor for each file type.
"""

import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
import extract_msg
import easyocr
import os
from src.common.interfaces import BaseExtractor
from src.common.image_classifier import is_scanned_document

# --- Global OCR Singleton ---
# We use a global variable so the model loads once per worker process, 
# rather than reloading for every single image file.
ocr_reader = None

def get_ocr_reader():
    global ocr_reader
    if ocr_reader is None:
        print(f"👁️  [Process {os.getpid()}] Loading OCR Model (GPU)...")
        ocr_reader = easyocr.Reader(['en'], gpu=True)
    return ocr_reader

# --- 1. PDF Extractor ---
class PDFExtractor(BaseExtractor):
    def extract(self, file_path: str):
        try:
            with fitz.open(file_path) as doc:
                for page_num, page in enumerate(doc, start=1):
                    text = page.get_text()
                    if text.strip():
                        yield (page_num, text)
        except Exception as e:
            print(f"⚠️ PDF Error {file_path}: {e}")

# --- 2. Word/Docx Extractor ---
class DocxExtractor(BaseExtractor):
    def extract(self, file_path: str):
        try:
            doc = docx.Document(file_path)
            # Word docs don't have strict pages, so we treat it as Page 1
            full_text = "\n".join([p.text for p in doc.paragraphs])
            if full_text.strip():
                yield (1, full_text)
        except Exception as e:
            print(f"⚠️ Docx Error {file_path}: {e}")

# --- 3. Spreadsheet Extractor (Head + Tail) ---
class SpreadsheetExtractor(BaseExtractor):
    def extract(self, file_path: str):
        try:
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            
            # Remove empty rows
            df.dropna(how='all', inplace=True)
            
            # 1. Capture Headers (Top 5 rows)
            head = df.head(5).to_string(index=False)
            
            # 2. Capture Totals (Bottom 5 rows)
            tail = ""
            if len(df) > 10:
                tail = "\n...[Middle Rows Skipped]...\n" + df.tail(5).to_string(index=False, header=False)
                
            full_text = f"--- START OF SHEET ---\n{head}\n{tail}\n--- END OF SHEET ---"
            yield (1, full_text)
            
        except Exception as e:
            print(f"⚠️ Excel Error {file_path}: {e}")

# --- 4. PowerPoint Extractor ---
class SlideExtractor(BaseExtractor):
    def extract(self, file_path: str):
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, start=1):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
                
                full_text = "\n".join(text_runs)
                if full_text.strip():
                    yield (i, full_text)
        except Exception as e:
            print(f"⚠️ PPTX Error {file_path}: {e}")

# --- 5. Email (.msg) Extractor ---
class EmailExtractor(BaseExtractor):
    def extract(self, file_path: str):
        try:
            msg = extract_msg.Message(file_path)
            content = (
                f"Subject: {msg.subject}\n"
                f"From: {msg.sender}\n"
                f"Date: {msg.date}\n"
                f"Body:\n{msg.body}"
            )
            msg.close()
            yield (1, content)
        except Exception as e:
            print(f"⚠️ Email Error {file_path}: {e}")

# --- 6. Image/OCR Extractor ---
class ImageExtractor(BaseExtractor):
    def extract(self, file_path: str):
        # 1. Heuristic Check: Is this a document?
        if is_scanned_document(file_path):
            try:
                reader = get_ocr_reader()
                # detail=0 returns a list of strings
                results = reader.readtext(file_path, detail=0)
                full_text = " ".join(results)
                
                if full_text.strip():
                    yield (1, full_text)
            except Exception as e:
                print(f"⚠️ OCR Error {file_path}: {e}")

# --- 7. Plain Text Extractor ---
class TextExtractor(BaseExtractor):
    def extract(self, file_path: str):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                yield (1, f.read())
        except Exception as e:
            print(f"⚠️ Text Error {file_path}: {e}")